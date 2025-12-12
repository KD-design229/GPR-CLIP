#!/usr/bin/env python3
"""
代码整体分析与模型架构总结 - PPT生成脚本
Generate comprehensive architecture analysis PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_title_slide(prs, title, subtitle=""):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide

def create_content_slide(prs, title, content_list):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for item in content_list:
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
            
        p = text_frame.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(18 - level * 2)
    
    return slide

def create_architecture_diagram_slide(prs, title, components):
    """创建架构图页面"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add boxes for each component
    top = Inches(1.5)
    left = Inches(1)
    width = Inches(2.5)
    height = Inches(1)
    spacing = Inches(0.3)
    
    for i, (comp_title, comp_desc) in enumerate(components):
        # Calculate position (max 3 per row)
        row = i // 3
        col = i % 3
        
        box_left = left + col * (width + spacing)
        box_top = top + row * (height + spacing)
        
        # Add rectangle
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            box_left, box_top, width, height
        )
        
        # Set fill color based on index
        colors = [
            RGBColor(91, 155, 213),   # Blue
            RGBColor(237, 125, 49),   # Orange
            RGBColor(165, 165, 165),  # Gray
            RGBColor(255, 192, 0),    # Yellow
            RGBColor(112, 173, 71),   # Green
            RGBColor(158, 72, 178),   # Purple
        ]
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = colors[i % len(colors)]
        
        # Add text
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = comp_title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Add description
        p2 = text_frame.add_paragraph()
        p2.text = comp_desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def generate_ppt():
    """生成完整的PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: 封面
    create_title_slide(
        prs,
        "FedDWA 联邦学习框架",
        "代码架构分析与模型总结\nPersonalized Federated Learning with Dynamic Weight Adjustment"
    )
    
    # Slide 2: 项目概述
    create_content_slide(prs, "项目概述 (Project Overview)", [
        "📚 研究来源: IJCAI 2023 论文",
        "🎯 核心目标: 个性化联邦学习 (Personalized Federated Learning)",
        "🔬 主要创新: 动态权重调整 (Dynamic Weight Adjustment)",
        "",
        "✨ 支持的联邦学习算法:",
        (1, "FedDWA - 动态权重聚合 (本项目核心)"),
        (1, "FedAvg - 联邦平均 (基线方法)"),
        (1, "FedProx - 近端正则化"),
        (1, "FedNova - 归一化平均"),
        (1, "FedSAM - 锐度感知最小化"),
        (1, "MOON - 模型对比学习"),
    ])
    
    # Slide 3: 系统架构
    create_content_slide(prs, "整体系统架构", [
        "🏗️ 三层架构设计:",
        "",
        "1️⃣ Server Layer (服务器层)",
        (1, "负责客户端选择、模型聚合、全局协调"),
        (1, "ServerBase: 基础服务器类"),
        (1, "FedDWA/FedAvg/FedProx等: 算法特定实现"),
        "",
        "2️⃣ Client Layer (客户端层)",
        (1, "执行本地训练、模型更新"),
        (1, "ClientBase: 基础客户端类"),
        (1, "ClientFedDWA: 支持两步模型预测"),
        "",
        "3️⃣ Model Layer (模型层)",
        (1, "多种神经网络架构"),
        (1, "从简单CNN到Transformer全覆盖"),
    ])
    
    # Slide 4: FedDWA核心算法
    create_content_slide(prs, "FedDWA 核心算法", [
        "💡 动态权重调整 (Dynamic Weight Adjustment)",
        "",
        "🔹 核心思想:",
        (1, "为每个客户端计算个性化的聚合权重"),
        (1, "基于模型相似度选择最相关的邻居"),
        (1, "Top-K选择机制保留最有用的模型"),
        "",
        "🔹 算法流程:",
        (1, "1. 客户端本地训练 → 生成当前模型 w_t"),
        (1, "2. 额外训练一步 → 生成下一步模型 w_{t+1}"),
        (1, "3. 服务器计算权重矩阵: W[i,j] ∝ 1/||w_{j,t+1} - w_{i,t}||²"),
        (1, "4. 列归一化 + Top-K剪枝 → 最优权重矩阵"),
        (1, "5. 个性化聚合: w_{i,new} = Σ W[j,i] * w_{j,t}"),
    ])
    
    # Slide 5: 支持的模型架构 (1)
    create_content_slide(prs, "支持的模型架构 (1/3) - 基础模型", [
        "🧱 基础卷积神经网络:",
        "",
        "1. CIFAR10Model / CIFAR100Model",
        (1, "架构: Conv → BN → ReLU → MaxPool × 2 + FC"),
        (1, "参数: 2304 → 512 → 10/100 classes"),
        (1, "特点: 支持Head/Body分离 (用于个性化)"),
        "",
        "2. FedAvgCNN",
        (1, "经典联邦学习基线模型"),
        (1, "5×5卷积核 + 2×2最大池化"),
        "",
        "3. ResNet8 / ResNet18",
        (1, "残差网络变体"),
        (1, "支持CIFAR-10/100, Tiny-ImageNet, GPR"),
    ])
    
    # Slide 6: 支持的模型架构 (2)
    create_content_slide(prs, "支持的模型架构 (2/3) - 现代架构", [
        "🚀 现代高效架构:",
        "",
        "1. MobileViT (Vision Transformer for Mobile)",
        (1, "使用 timm.mobilevit_s 作为backbone"),
        (1, "GPR模式: 添加可学习的信号预处理层"),
        (1, "特点: 轻量级、适合移动端部署"),
        "",
        "2. EfficientNet-B0",
        (1, "高效卷积网络 (来自 timm 库)"),
        (1, "复合缩放策略优化深度/宽度/分辨率"),
        "",
        "3. ResNet18-TIMM",
        (1, "标准化的 ResNet18 实现"),
        (1, "支持多数据集 (CIFAR/Tiny-ImageNet/GPR)"),
    ])
    
    # Slide 7: 支持的模型架构 (3)
    create_content_slide(prs, "支持的模型架构 (3/3) - 前沿模型", [
        "🌟 前沿多模态与专用架构:",
        "",
        "1. FedCLIP (CLIP for Federated Learning)",
        (1, "基于OpenAI CLIP的多模态模型"),
        (1, "核心组件: MaskedMLP + Softmax注意力机制"),
        (1, "支持CoOp (Context Optimization)"),
        (1, "物理先验初始化: GPR B-scan signal..."),
        (1, "Prompt Ensemble: 多模板聚合提升鲁棒性"),
        "",
        "2. GPR-FedSense (Ground Penetrating Radar)",
        (1, "专为探地雷达数据设计"),
        (1, "三层架构: 本地私有层 + 全局共享层 + 个性化头"),
        (1, "支持FedVLS和FedDecorr优化策略"),
    ])
    
    # Slide 8: FedCLIP详细架构
    create_content_slide(prs, "FedCLIP 架构详解", [
        "🖼️ 视觉-语言联邦学习模型",
        "",
        "核心组件:",
        "1. CLIP Backbone (冻结)",
        (1, "使用预训练的ViT-B/32或ViT-L/14"),
        (1, "图像编码器: encode_image() → 512D特征"),
        (1, "文本编码器: encode_text() → 512D特征"),
        "",
        "2. MaskedMLP Adapter (可训练)",
        (1, "二值化步长函数: 动态生成稀疏掩码"),
        (1, "可学习阈值: 控制剪枝力度"),
        (1, "Attention机制: fea_attn = MaskedMLP → BN → ReLU → MaskedMLP → Softmax"),
        "",
        "3. CoOp Prompt Learning (可选)",
        (1, "PromptLearner: 学习 n_ctx=16 个上下文向量"),
        (1, "Physics-based初始化: 使用GPR领域先验"),
    ])
    
    # Slide 9: GPR-FedSense详细架构
    create_content_slide(prs, "GPR-FedSense 架构详解", [
        "📡 探地雷达专用联邦学习模型",
        "",
        "三层分离式设计:",
        "",
        "1️⃣ 本地私有层 (不参与聚合)",
        (1, "GPRSignalNorm: 可学习的信号归一化"),
        (1, "GPRFeatureExtractor: 时空特征提取"),
        (2, "时间域卷积 (5×1): 捕获深度反射"),
        (2, "空间域卷积 (1×5): 捕获横向延续性"),
        "",
        "2️⃣ 全局共享层 (联邦聚合)",
        (1, "支持三种Backbone: CNN / ResNet18 / MobileViT"),
        (1, "提取高层语义特征 → 512D"),
        "",
        "3️⃣ 个性化分类头 (ALA自适应聚合)",
        (1, "Dropout → FC(512→256) → FC(256→8类)"),
        (1, "支持ALA加权聚合应对Non-IID"),
    ])
    
    # Slide 10: 客户端-服务器交互流程
    create_content_slide(prs, "客户端-服务器交互流程", [
        "🔄 联邦学习训练循环 (每轮):",
        "",
        "1. 服务器选择客户端",
        (1, "随机选择 client_frac × N 个客户端"),
        (1, "例: N=20, frac=0.5 → 10个客户端参与"),
        "",
        "2. 服务器分发模型",
        (1, "send_models(): 发送全局模型/个性化模型"),
        "",
        "3. 客户端本地训练",
        (1, "本地数据集上训练 E 个epoch"),
        (1, "FedDWA: 额外计算下一步模型"),
        "",
        "4. 客户端上传模型",
        (1, "receive_models(): 收集更新后的模型"),
        "",
        "5. 服务器聚合",
        (1, "FedAvg: 加权平均"),
        (1, "FedDWA: 动态权重矩阵聚合"),
    ])
    
    # Slide 11: 数据处理与Non-IID设置
    create_content_slide(prs, "数据处理与Non-IID设置", [
        "📊 数据集支持:",
        (1, "CIFAR-10 / CIFAR-100"),
        (1, "CINIC-10"),
        (1, "Tiny-ImageNet (200类)"),
        (1, "GPR Custom (探地雷达 8类)"),
        "",
        "🔀 Non-IID分布类型:",
        "",
        "Type 8 - 病态异构 (Pathological)",
        (1, "每个客户端只有少数几个类别"),
        "",
        "Type 9 - 实际异构1 (Dirichlet分布)",
        (1, "使用Dirichlet(α) 控制数据分布偏斜"),
        (1, "α越小，Non-IID程度越高"),
        "",
        "Type 10 - 实际异构2 (类别数+比例)",
        (1, "每个客户端固定num_types个类,占比ratio"),
    ])
    
    # Slide 12: 优化策略
    create_content_slide(prs, "高级优化策略", [
        "🚀 支持的优化技术:",
        "",
        "1. FedVLS (Vacant-class Distillation)",
        (1, "空置类蒸馏: 处理本地缺失的类别"),
        (1, "使用全局模型作为教师模型"),
        "",
        "2. FedDecorr (Feature Decorrelation)",
        (1, "特征去相关: 减少特征冗余"),
        (1, "正交约束提升泛化能力"),
        "",
        "3. ALA (Adaptive Layer Aggregation)",
        (1, "自适应层聚合: 不同层使用不同权重"),
        (1, "参数: rand_percent, layer_idx, eta"),
        "",
        "4. Learning Rate Decay",
        (1, "学习率衰减: lr_decay, lr_decay_step"),
    ])
    
    # Slide 13: 代码模块结构
    create_architecture_diagram_slide(prs, "代码模块结构", [
        ("main.py", "入口程序\n参数解析"),
        ("servers/", "服务器实现\nBase+算法"),
        ("clients/", "客户端实现\nBase+算法"),
        ("model/", "神经网络\nMLModel.py"),
        ("utils/", "工具函数\n数据/日志"),
        ("logs_feddwa/", "实验日志\nCSV结果"),
    ])
    
    # Slide 14: 实验配置
    create_content_slide(prs, "实验配置参数", [
        "⚙️ 主要超参数:",
        "",
        "全局参数:",
        (1, "Tg: 全局通信轮数 (默认100)"),
        (1, "client_num: 客户端总数 (默认20)"),
        (1, "client_frac: 每轮参与比例 (默认0.5)"),
        "",
        "本地训练:",
        (1, "E: 本地训练epoch数 (默认1)"),
        (1, "B: 本地batch size (默认20)"),
        (1, "lr: 学习率 (默认0.01)"),
        "",
        "FedDWA特定:",
        (1, "feddwa_topk: Top-K邻居数 (默认5)"),
        (1, "next_round: 下一步预测轮数 (默认1)"),
    ])
    
    # Slide 15: 结果保存与可视化
    create_content_slide(prs, "结果保存与可视化", [
        "📈 自动生成的结果:",
        "",
        "1. 训练日志 (logs_feddwa/)",
        (1, "*.json: 测试准确率、训练损失等"),
        (1, "*.csv: 详细的逐轮结果"),
        (1, "*_model_structure.json: 模型结构"),
        "",
        "2. 混淆矩阵 (Confusion Matrices)",
        (1, "client_confusion_matrices/: 每个客户端的CM"),
        (1, "使用seaborn绘制热力图"),
        "",
        "3. 训练曲线 (plot_training_results)",
        (1, "准确率 vs. 轮数"),
        (1, "损失 vs. 轮数"),
        "",
        "4. 模型检查点",
        (1, "保存最优模型权重"),
    ])
    
    # Slide 16: 关键创新点总结
    create_content_slide(prs, "关键创新点总结", [
        "💡 本项目的核心贡献:",
        "",
        "1️⃣ FedDWA算法",
        (1, "基于模型相似度的动态权重聚合"),
        (1, "为每个客户端提供个性化模型"),
        (1, "Top-K机制提高聚合效率"),
        "",
        "2️⃣ 多模态联邦学习 (FedCLIP)",
        (1, "首次将CLIP引入联邦学习"),
        (1, "CoOp物理先验初始化"),
        (1, "MaskedMLP实现稀疏适配"),
        "",
        "3️⃣ 专用领域适配 (GPR-FedSense)",
        (1, "三层分离架构适配探地雷达"),
        (1, "时空特征提取器"),
        (1, "本地私有层保护设备特异性"),
        "",
        "4️⃣ 全面的算法对比框架",
        (1, "统一接口支持6种联邦学习算法"),
    ])
    
    # Slide 17: 技术栈
    create_content_slide(prs, "技术栈 (Tech Stack)", [
        "🛠️ 主要依赖库:",
        "",
        "深度学习框架:",
        (1, "PyTorch - 核心训练框架"),
        (1, "torchvision - 视觉模型与数据集"),
        (1, "timm - 预训练模型库"),
        (1, "CLIP - 多模态预训练模型"),
        "",
        "科学计算:",
        (1, "NumPy - 数值计算"),
        (1, "scikit-learn - 评估指标"),
        "",
        "可视化:",
        (1, "matplotlib - 绘图"),
        (1, "seaborn - 统计可视化"),
        "",
        "其他:",
        (1, "h5py - HDF5数据格式"),
        (1, "pathlib - 路径管理"),
    ])
    
    # Slide 18: 模型性能对比
    create_content_slide(prs, "模型复杂度对比", [
        "📊 不同模型的参数规模:",
        "",
        "轻量级模型:",
        (1, "CIFAR10Model: ~2.3M 参数"),
        (1, "FedAvgCNN: ~1.2M 参数"),
        "",
        "中等模型:",
        (1, "ResNet18: ~11M 参数"),
        (1, "MobileViT-S: ~5M 参数"),
        (1, "EfficientNet-B0: ~5M 参数"),
        "",
        "大型模型:",
        (1, "FedCLIP (ViT-B/32): ~87M 参数 (冻结)"),
        (1, "  + Adapter: ~0.5M 可训练参数"),
        "",
        "专用模型:",
        (1, "GPR-FedSense: 可配置 (3M~15M)"),
    ])
    
    # Slide 19: 应用场景
    create_content_slide(prs, "应用场景 (Application Scenarios)", [
        "🌐 联邦学习适用领域:",
        "",
        "1. 医疗健康 (Healthcare)",
        (1, "医院间协作训练，无需共享病患数据"),
        (1, "FedCLIP: 医学图像+报告联合学习"),
        "",
        "2. 智能交通 (Intelligent Transportation)",
        (1, "车辆间协同感知"),
        (1, "GPR-FedSense: 路面检测"),
        "",
        "3. 金融风控 (Finance)",
        (1, "银行间反欺诈模型"),
        (1, "保护客户隐私"),
        "",
        "4. 工业检测 (Industrial Inspection)",
        (1, "探地雷达、无损检测"),
        (1, "设备间知识共享"),
    ])
    
    # Slide 20: 未来展望
    create_content_slide(prs, "未来工作方向", [
        "🔮 可能的改进方向:",
        "",
        "1. 算法优化",
        (1, "结合差分隐私 (Differential Privacy)"),
        (1, "安全多方计算 (Secure Multi-party Computation)"),
        (1, "拜占庭鲁棒聚合 (Byzantine-robust Aggregation)"),
        "",
        "2. 模型压缩",
        (1, "知识蒸馏 (Knowledge Distillation)"),
        (1, "模型剪枝 (Pruning)"),
        (1, "量化感知训练 (Quantization-aware Training)"),
        "",
        "3. 通信优化",
        (1, "梯度压缩 (Gradient Compression)"),
        (1, "部分参数更新 (Partial Update)"),
        "",
        "4. 异步联邦学习",
        (1, "支持异构设备异步训练"),
    ])
    
    # Slide 21: 结论
    create_content_slide(prs, "总结 (Conclusion)", [
        "✅ 本项目实现了:",
        "",
        "🏆 完整的联邦学习框架",
        (1, "支持6种主流联邦学习算法"),
        (1, "统一的Server-Client架构"),
        "",
        "🧠 丰富的模型库",
        (1, "从基础CNN到前沿Transformer"),
        (1, "特别优化的GPR/CLIP模型"),
        "",
        "🔬 严谨的实验设计",
        (1, "多种Non-IID设置"),
        (1, "完善的评估与可视化"),
        "",
        "📚 清晰的代码组织",
        (1, "模块化设计，易于扩展"),
        (1, "详细的注释与文档"),
        "",
        "💬 Thank you for your attention!",
    ])
    
    # Save presentation
    output_path = "/home/engine/project/FedDWA_Architecture_Analysis.pptx"
    prs.save(output_path)
    print(f"✅ PPT生成成功! 保存路径: {output_path}")
    return output_path

if __name__ == "__main__":
    generate_ppt()
